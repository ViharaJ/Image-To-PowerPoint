"""
Places images into an empty PowerPoint Presentation and creates a small table
at the bottom of each image.

It is assumed that images HAVE NOT been RENAMED. The excel file will be used to
figure out the experiment run, conditions, etc and eventually placement on the 
slide.


How to use: 
    1. Create an empty presentation. Change the variable 'path' to its location 
    2. Change 'imagePath' to folder with seite and partikel images
    3. Change 'hellDunkelPath' 
    4. Change excelPath to excel file containing all the data 
    5. Run 
    
IMPORTANT: This script is not very versatile and needs to be refactored for 
a more general use case. However, this script can be used as a quick start guide
for more general scripts. 

Documenation to pptx: https://python-pptx.readthedocs.io/
"""
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd 
import os 
import sys
from PIL import Image
import io
from pptx.enum.text import MSO_AUTO_SIZE
import timeit


def canPlaceSP(imgName, idx):
    imgName = imgName[:-4]
    comp = imgName.split('_')
    
    if(len(comp) < 4):
        return False
    
    if idx == 0 and comp[-1].lower() == "seite":
        return True
    elif idx == 1 and comp[-1].lower() == "partikel":
        return True
    
    return False


def canPlaceHD(imgName, idx):
    imgName = imgName[:-4]
    comp = imgName.split('_')
    
    if(len(comp) < 4):
        return False
    
    if idx == 0 and comp[-1].lower() == "dunkel":
        return True
    elif idx == 1 and comp[-1].lower() == "hell":
        return True
    
    return False


def addToDict(directory, Dict):
    for filename in directory: 
        #Check if file is of accepted type
        if( '.' in filename and filename.split('.')[-1].lower() in acceptedFileTypes):
            nameComp = filename[:-4].split('_')
            if(len(nameComp) > 1):
                Versuch = nameComp[0]
                Run = nameComp[1]
            
                val = Dict.get(Versuch + '_' + Run)
                
                if val:
                    val.append(filename)
                    Dict.update({Versuch + '_' + Run: val})
                else:
                    Dict[Versuch + '_' + Run] = [filename]
                    
    
                    

start = timeit.default_timer()
    
path = "C:/Users/v.jayaweera/Documents/Tim/Slides/TestSlideEmptyC2.pptx"
imagePath = "C:/Users/v.jayaweera/Documents/Tim/Slides/20230607_Proben"
hellDunkelPath = "C:/Users/v.jayaweera/Documents/Tim/Slides/20230607_Proben im Pulverbett"
excelPath = "C:/Users/v.jayaweera/Documents/Tim/Microscope/20230607_Versuchsplan.xlsx"
acceptedFileTypes = ["jpg", "png", "bmp", "tif", "JPG"]

dirPictures = os.listdir(imagePath)
hdPictures = os.listdir(hellDunkelPath)
kvDict = dict()

#Add image file names to dictionary
addToDict(dirPictures, kvDict)
addToDict(hdPictures, kvDict)

            
#Open excel doc and presentation
df = pd.read_excel(excelPath)
prs = Presentation(path)


#Iterate over slides, 2 at a time, 
for i in range(1,len(prs.slides)-1,2):
    slide1 = prs.slides[i]
    slide2 = prs.slides[i+1]
    title = slide1.shapes.title.text
    
    subDf = df.loc[df['Simulate'] == title]
    #Sort sub data frame to ensure Runs are in ascending order
    subDF = subDf.sort_values('Run')
    
    #Get unique Runs, should return 3 items
    uniquePrefixes = (subDf.Versuch.astype(str) + '_' + subDf.Run.astype(str)).unique()
    
    #Slide 1
    left = 1
    Height = 6.38 
    Width = 10
    spacing = 1
    
    
    for j in range(len(uniquePrefixes)):
        #Images available to insert into slide
        imgSet = kvDict.get(uniquePrefixes[j])
        
              
        #if image set exists, insert into Slide 1
        if imgSet: 
            top = 3.6           
            for j in range(2):
                imgIdx = 0
                while(imgIdx < len(imgSet)):
                    currImage = imgSet[imgIdx]
                    
                    if(canPlaceSP(currImage, j)):
                        #Reduce image size
                        img = Image.open(imagePath + '/'+currImage)
                        img.thumbnail((500, 500), Image.LANCZOS)
                        image_stream = io.BytesIO()
                        img.save(image_stream, "PNG")
                        image_stream.seek(0)
                        
                        slide1.shapes.add_picture(image_stream, Cm(left), Cm(top), height=Cm(Height))
                        break
                    else:
                        imgIdx = imgIdx + 1 
                        
                    
                top = top + Height + 0.25
               
                
        left = left + Width + spacing
        
    #add table to slide 1
    tableW = Width*3 + spacing*3
    tableHeight = 1
    shape = slide1.shapes.add_table(1, 11, Cm(1), Cm(top-0.1), Cm(tableW),Cm(tableHeight))  
    table = shape.table
    
    tIndx = 0
    for j in range(3):
        row = subDf.iloc[j]
        
        for k in range(3):            
            table.cell(0,tIndx + k).text = row["P [W]"].astype(str) + "W, " + row["Vs [mm/s]"].astype(str) + "mm/s"
            table.cell(0,tIndx + k).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            table.columns[tIndx + k].width = Inches(1.25)
       
        if(j < 2):
            table.columns[tIndx + 3].width = Cm(1.5)
            
        tIndx = tIndx + 4
        
            
    #Slide 2, hell und dunkel
    left = 1
    for j in range(len(uniquePrefixes)):
        #Images available to insert into slide
        imgSet = kvDict.get(uniquePrefixes[j])
        
              
        #if image set exists, insert
        if imgSet: 
            top = 3.6           
            for j in range(2):
                imgIdx = 0
                while(imgIdx < len(imgSet)):
                    currImage = imgSet[imgIdx]
                    
                    if(canPlaceHD(currImage, j)):
                        #Reduce image size
                        img = Image.open(hellDunkelPath + '/'+currImage)
                        img.thumbnail((500, 500), Image.LANCZOS)
                        image_stream = io.BytesIO()
                        img.save(image_stream, "PNG")
                        image_stream.seek(0)
                        
                        slide2.shapes.add_picture(image_stream, Cm(left), Cm(top), height=Cm(Height))
                        break
                    else:
                        imgIdx = imgIdx + 1 
                        
                    
                top = top + Height + 0.25
               

        left = left + Width + spacing
        
    #add table to slide2
    shape = slide2.shapes.add_table(1, 11, Cm(1), Cm(top-0.1) ,Cm(tableW),Cm(tableHeight))  
    table = shape.table
    
    tIndx = 0
    for j in range(3):
        row = subDf.iloc[j]
        
        for k in range(3):            
            table.cell(0,tIndx + k).text = row["P [W]"].astype(str) + "W, " + row["Vs [mm/s]"].astype(str) + "mm/s"
            table.cell(0,tIndx + k).text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            table.columns[tIndx + k].width = Inches(1.25)
       
        if(j < 2):
            table.columns[tIndx + 3].width = Cm(1.5)
            
        tIndx = tIndx + 4
        
prs.save(path)

stop = timeit.default_timer()

print('Time: ', stop - start)